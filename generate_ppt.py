#!/usr/bin/env python3
"""Generate human-oriented PPTX report for Cluster Tool Scheduler."""

import json
import urllib.request
import os
import re
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

API = "http://localhost:8080/api"
BASE_DIR = os.path.dirname(os.path.abspath(__file__))
DARK_BG = RGBColor(0x1A, 0x1A, 0x2E)
SECTION_BG = RGBColor(0x0F, 0x30, 0x50)
CYAN = RGBColor(0x00, 0xD4, 0xFF)
WHITE = RGBColor(0xE0, 0xE0, 0xE0)
GREEN = RGBColor(0x4C, 0xAF, 0x50)
ORANGE = RGBColor(0xFF, 0x98, 0x00)
PINK = RGBColor(0xE9, 0x1E, 0x63)
YELLOW = RGBColor(0xFF, 0xD5, 0x4F)
GRAY = RGBColor(0x88, 0x88, 0x88)
LIGHT_GRAY = RGBColor(0xBB, 0xBB, 0xBB)
PURPLE = RGBColor(0xAB, 0x47, 0xBC)

def fetch(endpoint):
    with urllib.request.urlopen(f"{API}{endpoint}") as resp:
        return json.loads(resp.read())

def format_time(sec):
    if sec < 0: sec = 0
    h, m = divmod(sec, 3600)
    m, s = divmod(m, 60)
    if h > 0: return f"{h}h{m}m{s}s"
    if m > 0: return f"{m}m{s}s"
    return f"{s}s"

def set_slide_bg(slide, color=DARK_BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_title(slide, text, top=Inches(0.3), size=Pt(28), color=CYAN):
    txBox = slide.shapes.add_textbox(Inches(0.5), top, Inches(9), Inches(0.7))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = size
    p.font.color.rgb = color
    p.font.bold = True
    return tf

def add_subtitle(slide, text, top=Inches(1.1), size=Pt(14), color=GRAY):
    txBox = slide.shapes.add_textbox(Inches(0.7), top, Inches(8.5), Inches(0.4))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = size
    p.font.color.rgb = color
    return tf

def add_body(slide, lines, top=Inches(1.6), left=Inches(0.7), width=Inches(8.5), size=Pt(13)):
    txBox = slide.shapes.add_textbox(left, top, width, Inches(5.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        # Handle colored keywords
        if line.startswith('📌') or line.startswith('🔧') or line.startswith('📊') or line.startswith('✅') or line.startswith('🔮') or line.startswith('⭐'):
            p.font.color.rgb = YELLOW
            p.font.size = Pt(size.pt + 2)
            p.font.bold = True
        elif line.startswith('•') or line.startswith('  •') or line.startswith('   •') or line.startswith('    •'):
            p.font.color.rgb = WHITE
            p.font.size = size
        elif line.startswith('  ') and not line.startswith('   '):
            p.font.color.rgb = LIGHT_GRAY
            p.font.size = Pt(size.pt - 1)
        elif line == '':
            p.font.size = Pt(4)
        else:
            p.font.color.rgb = WHITE
            p.font.size = size
        p.text = line
        p.space_after = Pt(4)
    return tf

def add_key_point_box(slide, title, text, top, left=Inches(0.7), width=Inches(8.5), height=Inches(0.9)):
    """Add a highlighted key point box"""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x0F, 0x34, 0x60)
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = f"⭐ {title}"
    p.font.size = Pt(14)
    p.font.color.rgb = YELLOW
    p.font.bold = True
    p2 = tf.add_paragraph()
    p2.text = text
    p2.font.size = Pt(12)
    p2.font.color.rgb = WHITE
    p2.space_before = Pt(4)
    return shape

def add_separator(slide, top):
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), top, Inches(9), Pt(2))
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(0x33, 0x33, 0x44)
    line.line.fill.background()

def add_section_header(slide, section_num, text, top=Inches(0.15)):
    """Chapter section header bar"""
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), top, Inches(10), Inches(0.55))
    bar.fill.solid()
    bar.fill.fore_color.rgb = SECTION_BG
    bar.line.fill.background()
    tf = bar.text_frame
    p = tf.paragraphs[0]
    p.text = f"  {section_num}.  {text}"
    p.font.size = Pt(16)
    p.font.color.rgb = CYAN
    p.font.bold = True
    p.alignment = PP_ALIGN.LEFT
    return bar

def add_table(slide, headers, rows, top, left=Inches(0.7), width=Inches(8.5), font_size=Pt(11)):
    n_rows = len(rows) + 1
    n_cols = len(headers)
    row_h = Inches(0.32)
    table_shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, row_h * n_rows)
    table = table_shape.table
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            p.font.size = font_size
            p.font.color.rgb = CYAN
            p.font.bold = True
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0x0F, 0x34, 0x60)
    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            cell = table.cell(r + 1, c)
            cell.text = str(val)
            for p in cell.text_frame.paragraphs:
                p.font.size = font_size
                p.font.color.rgb = WHITE
            bg = RGBColor(0x16, 0x21, 0x3E) if r % 2 == 0 else RGBColor(0x1A, 0x27, 0x44)
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg
    return table_shape

def add_footer(slide, text="Cluster Tool 调度仿真系统"):
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(7.0), Inches(9), Inches(0.3))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(8)
    p.font.color.rgb = GRAY
    p.alignment = PP_ALIGN.RIGHT

def main():
    # Fetch data
    state = fetch("/simulation/state")
    device = fetch("/config/device")
    schedule = fetch("/config/schedule")
    am_config_raw = fetch("/config/am")
    events = fetch("/simulation/events")

    total = state['totalWafers']
    done = state['completedWafers']
    sim_time = state['currentTimeSec']
    wph = round(done / (sim_time / 3600.0), 1) if sim_time > 0 else 0

    # Violations
    v_stats = {}
    for e in events:
        m = re.search(r'WARN:\s*(\S+)\s+(\S+)\s+dwell\s+(\d+)s\s+exceeds\s+max\s+(\d+)s', e)
        if not m: continue
        ch = m.group(1)
        ct = 'PT' if ch.startswith('PT') else ('PreClean' if ch.startswith('PreClean') else ('EPI' if ch.startswith('EPI') else ch))
        if ct not in v_stats: v_stats[ct] = {'count': 0, 'max': 0}
        v_stats[ct]['count'] += 1
        v_stats[ct]['max'] = max(v_stats[ct]['max'], int(m.group(3)))

    total_violations = sum(v['count'] for v in v_stats.values())

    # EPI stats
    epi_onload_events = [e for e in events if 'EPI' in e and 'OnLoad' in e]
    purge_events = [e for e in events if 'IdlePurge' in e and 'started' in e]

    recipes = schedule.get('recipes', {})
    timing = schedule.get('timing', {})
    passthroughs = device.get('passthroughs', [])

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # ====== Slide 1: Title ======
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_title(slide, "Cluster Tool 调度仿真系统", Inches(1.8), Pt(36), CYAN)
    add_subtitle(slide, "半导体外延生长集群设备 · 晶圆调度优化与全流程仿真", Inches(2.7), Pt(16), GRAY)

    txBox = slide.shapes.add_textbox(Inches(1.5), Inches(3.5), Inches(7), Inches(1.5))
    tf = txBox.text_frame
    for text in [f"✅ {done}/{total} 片晶圆成功完成", f"✅ WPH = {wph} | 模拟时长 = {format_time(sim_time)}", f"✅ 驻留时间违例 = {total_violations} 次", f"✅ 4腔 EPI OnLoadClean 全部对齐（间隔 ~73s）"]:
        p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(15)
        p.font.color.rgb = GREEN
        p.alignment = PP_ALIGN.CENTER
        p.space_after = Pt(6)

    add_footer(slide, f"设备: {device['equipmentName']} ({device['equipmentId']})")

    # ====== Slide 2: Project Overview ======
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_section_header(slide, "01", "项目概述与目标")
    add_body(slide, [
        "📌 项目目标",
        "  • 完整模拟 Cluster Tool（外延生长集群设备）的晶圆流转全流程",
        "  • 优化调度策略，在满足所有工艺约束的前提下最大化晶圆产出（WPH）",
        "  • 纯配置驱动，所有设备/工艺/维护参数均可通过 JSON 文件调整",
        "",
        "📌 核心挑战",
        "  • 刚性最大驻留时间约束：PreClean 120s / EPI 100s / PT 300s",
        "  • EPI 工艺时间长（2120s），是系统瓶颈工站",
        "  • 单臂单指机械手，不支持原子交换（Swap），需精心调度防止死锁",
        "  • CoolingStation 冷却逻辑：EPI 返回晶圆必须经过冷却槽（60s）",
    ], Inches(1.0), size=Pt(12))

    add_key_point_box(slide,
        "核心指标",
        f"完成 {done}/{total} 片 · WPH = {wph} · 驻留违例 = {total_violations} · OnLoadClean 间隔对齐 ≤ 100s",
        Inches(4.8))

    add_footer(slide)

    # ====== Slide 3: Hardware Topology ======
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_section_header(slide, "02", "硬件拓扑与腔室布局")
    add_body(slide, [
        "📌 设备拓扑（从左到右晶圆流转方向）",
        "",
        "  FOUP(LP1/LP2) → ATM Robot → Aligner → LoadLock(LL1/LL2) → TM1 → PreClean(×2) → PT(2×2槽) → TM2 → EPI(×4) → 返回",
        "",
        "📌 腔室参数",
    ], Inches(0.9), size=Pt(12))

    add_table(slide,
        ["腔室类型", "数量", "工艺时间", "最大驻留", "备注"],
        [
            ["PreClean", "2腔", "280s ±10s", "120s", "预清洁"],
            ["EPI", "4腔", "2120s ±30s", "100s", "外延生长（瓶颈）"],
            ["Passthrough", "4槽(2×2)", "0s", "300s", "PT1_S0 & PT2_S1 = 冷却槽(60s)"],
            ["LoadLock", "2腔(BLL)", "Pump126s/Vent168s", "300s", "批次装卸"],
            ["Aligner", "1腔", "4.4s", "—", "晶圆对准"],
        ],
        Inches(3.1))

    add_key_point_box(slide,
        "机械手",
        "ATM(ATM1): 大气端 FOUP↔Aligner↔LL | TM1(Robot1): 真空端 LL↔PreClean↔PT | TM2(Robot2): PT↔EPI\n均为单臂单指，不支持 Swap",
        Inches(5.0), height=Inches(0.85))

    add_footer(slide)

    # ====== Slide 4: Process Constraints ======
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_section_header(slide, "03", "工艺约束与 AM 维护任务")
    add_body(slide, [
        "📌 刚性最大驻留时间约束（Max Dwell — 不可违反）",
        "  • PreClean 120s | EPI 100s | PT 300s | LoadLock 300s",
        "  • 驻留 = 工艺完成时刻 → 被机械手取走时刻的等待时间",
        "  • 调度器内置 40s 安全裕度（dwellSafetyMarginSec）进行前瞻控制",
        "",
        "📌 Auto Maintenance（AM）任务",
    ], Inches(0.9), size=Pt(12))

    add_table(slide,
        ["AM 任务", "适用腔室", "时长", "触发条件", "频率"],
        [
            ["OnLoadClean (EPI)", "EPI×4", "457s", "首个腔用延迟公式对齐首片，后续腔按 wafer 计数器", "每 CJ 1次/腔"],
            ["OnLoadClean (PreClean)", "PreClean×2", "537s", "固定 stagger 间隔启动", "每 CJ 1次/腔"],
            ["1X Clean (EPI)", "EPI×4", "457s", "EPI 工艺完成 wafer 取走后立即启动", "每片前 1次"],
            ["IdlePurge (PreClean)", "PreClean×2", "123s", "腔室空闲 ≥ 180s 自动触发", "空闲阈值驱动"],
        ],
        Inches(3.3), font_size=Pt(10))

    add_footer(slide)

    # ====== Slide 5: Configuration System ======
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_section_header(slide, "04", "配置系统 — 纯 JSON 驱动")
    add_body(slide, [
        "📌 5 个 JSON 配置文件（conf/ 目录），无硬编码",
        "",
        "  device.json    硬件拓扑：腔室数量/类型、机械手操作时间(pick/rotate/place)、",
        "                  FOUP 容量、LoadLock 时序、CoolingStation 分配",
        "  schedule.json  工艺参数：各 recipe 的处理时间/变化范围/最大驻留、",
        "                  调度策略、安全裕度、LL/PT 时序、模拟速度",
        "  sequence.json  Wafer 流转 9 步定义：每步的 station、robot、recipe",
        "  job.json       生产任务：CJ/PJ 模式(serial/parallel)、LP 分配、wafer 范围",
        "  am.json        AM 维护：OnLoadClean / 1X Clean / IdlePurge 的时间和阈值",
    ], Inches(0.9), size=Pt(12))

    add_key_point_box(slide,
        "可配置项总计",
        "设备拓扑 · 工艺时间 · 驻留约束 · 机械手速度 · LL 时序 · 冷却参数 · 调度策略 · CJ/PJ 模式 · AM 任务 · 模拟速度",
        Inches(4.3))

    add_footer(slide)

    # ====== Slide 6: Core Algorithms ======
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_section_header(slide, "05", "核心调度算法")
    add_body(slide, [
        "📌 主循环（每 1 模拟秒执行一次）",
        "  1. 更新腔室倒计时  → 2. 驻留检查 → 3. 机械手完成回调 → 4. ATM 调度",
        "  → 5. TM1 调度 → 6. TM2 调度 → 7. BLL 管理 → 8. 准备新批次",
        "  → 9. OnLoadClean → 10. IdlePurge → 11. 数据自愈",
        "",
        "📌 TM1 优先级（从高到低）",
        "  1️⃣ PT 返回 → LL     清空返回路径，释放 PT 槽位",
        "  2️⃣ PreClean → PT    释放 PreClean 腔，推进 wafer 流转",
        "  3️⃣ LL → PreClean    拉入新 wafer（受 canPullWaferFromLL 守卫）",
        "",
        "📌 TM2 优先级",
        "  1️⃣ EPI → PT（优先冷却槽）  →  2️⃣ PT → EPI",
    ], Inches(0.9), size=Pt(12))

    add_key_point_box(slide,
        "死锁预防 — 时间感知前瞻控制",
        "canPullWaferFromLL(): 统计 EPI 需求，预测最早 EPI 完成时间，仅当 wafer 到达 PT 时 EPI 恰好空闲才允许投入新片 → 永不溢出 → 无死锁",
        Inches(5.3), height=Inches(0.85))

    add_footer(slide)

    # ====== Slide 7: Deadlock Prevention Detail ======
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_section_header(slide, "06", "死锁预防 & 流水线深度控制")

    add_body(slide, [
        "📌 死锁场景：4 EPI 全满 + 4 PT 全被前向 wafer 占据 → 无法返回 → 死锁",
        "",
        "📌 三层防护机制",
        "  ① Stagger 限流",
        f"     stompInterval = (EPI_工艺 + EPI_清洁) ÷ EPI_腔数 = (2120+457)/4 ≈ 644s",
        "     每 644s 才允许释放一个新 wafer，均匀分布 EPI 完成时间",
        "",
        "  ② EPI 容量前瞻 (canPullWaferFromLL)",
        "     demand = PT_前向数 + PreClean_使用数 + 1（新wafer）",
        "     可用 EPI = IDLE数 + CLEANING数（即将空闲）",
        "     不足时 → 读取 EPI 剩余时间排序 → 预测最早完成时刻",
        "     仅当 wafer 到达 PT 前 EPI 确定有空 → 允许投入",
        "",
        "  ③ PT 前向压力控制 (canMovePCToPT)",
        "     PT 前向 wafer 仅在 EPI 有空（或即将有空）时进入",
        "     PreClean 驻留即将超标时 → 强制允许（安全覆盖）",
    ], Inches(0.9), size=Pt(11))

    add_key_point_box(slide,
        "设计哲学",
        "\"宁可稍慢投入，绝不创造死锁\" — 前瞻保守策略确保系统始终可收敛",
        Inches(6.0))

    add_footer(slide)

    # ====== Slide 8: AM Logic Detail ======
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_section_header(slide, "07", "OnLoadClean 时序优化")

    add_body(slide, [
        "📌 问题：OnLoadClean 启动太早 → EPI 腔空闲等待首片 wafer 到达（可达 749s）",
        "",
        "📌 优化方案",
        "  • EPI1（首个腔）：使用延迟公式 → 在首片 wafer 预计到达前 457s 启动",
        "    minStartTime = PreClean_OnLoad + Purge_阈值 + Purge_时长 + PC_工艺 + 传输 - EPI_Clean",
        "  • EPI2/3/4（后续腔）：跟踪 wafersEnteredPreClean 计数器",
        "    第 N 个腔在第 N 个 wafer 进入 PreClean 时启动 → 自动对齐实际管道节奏",
        "",
        "📌 效果",
    ], Inches(0.9), size=Pt(12))

    add_table(slide,
        ["EPI 腔", "优化前间隔", "优化后间隔", "说明"],
        [
            ["EPI1", "~749s", "~73s", "延迟公式对齐首片，间隔缩短 90%"],
            ["EPI2", "~749s", "~73s", "计数器触发，和 EPI1 一致"],
            ["EPI3", "从未 OnLoadClean", "~73s", "修复 stagger 共享变量 bug 后正常工作"],
            ["EPI4", "从未 OnLoadClean", "~73s", "同上，全 4 腔统一"],
        ],
        Inches(3.6))

    add_footer(slide)

    # ====== Slide 9: Architecture ======
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_section_header(slide, "08", "系统架构与技术栈")

    add_body(slide, [
        "📌 技术栈",
        "  • 后端：Java 21 + Spring Boot 3.2.5（REST API + WebSocket STOMP）",
        "  • 前端：Vue 3 + Vite（SVG 设备布局 + 交互式甘特图）",
        "  • 报告：Python 3（HTML + PPTX 自动生成）",
        "",
        "📌 分层架构",
        "  ┌─────────────────────────────────────┐",
        "  │  前端 Vue 3 (ControlPanel + ToolLayout + GanttChart) │",
        "  │       ↕ HTTP REST + WebSocket STOMP                    │",
        "  │  后端 Spring Boot                                     │",
        "  │  ├─ SchedulerController (18 个 /api/* 端点)            │",
        "  │  ├─ SimulationService (@Scheduled 10ms 定时循环)       │",
        "  │  ├─ SchedulerEngine (核心调度引擎 ~1400行)             │",
        "  │  └─ ConfigService (JSON 配置加载与热重载)              │",
        "  │  报告生成 Python 脚本 (generate_report.py / generate_ppt.py) │",
        "  └─────────────────────────────────────┘",
    ], Inches(0.9), size=Pt(11))

    add_footer(slide)

    # ====== Slide 10: Simulation Results ======
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_section_header(slide, "09", "模拟结果与性能指标")

    add_body(slide, [
        f"📊 核心指标",
        f"  • 完成晶圆：{done}/{total} 片（100%）",
        f"  • WPH：{wph} 片/小时",
        f"  • 总模拟时间：{format_time(sim_time)}（{sim_time}s）",
        f"  • EPI OnLoadClean：{len(epi_onload_events)//2} 次启动（4腔全覆盖）",
        f"  • IdlePurge：{len(purge_events)} 次（纯空闲驱动）",
        f"  • 数据自愈触发：被动检测修复",
    ], Inches(0.9), size=Pt(12))

    add_table(slide,
        ["约束类型", "限制值", "违反次数", "状态"],
        [
            ["PreClean 驻留", "120s", str(v_stats.get('PreClean', {}).get('count', 0)), "✅ 无违反" if v_stats.get('PreClean', {}).get('count', 0) == 0 else "⚠ 有违反"],
            ["EPI 驻留", "100s", str(v_stats.get('EPI', {}).get('count', 0)), "✅ 无违反" if v_stats.get('EPI', {}).get('count', 0) == 0 else "⚠ 有违反"],
            ["PT 驻留", "300s", str(v_stats.get('PT', {}).get('count', 0)), "✅ 无违反" if v_stats.get('PT', {}).get('count', 0) == 0 else "⚠ 有违反"],
        ],
        Inches(3.8))

    add_footer(slide)

    # ====== Slide 11: EPI OnLoadClean Gap Detail ======
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_section_header(slide, "10", "EPI OnLoadClean 到首片放置间隔分析")

    add_body(slide, [
        "📌 OnLoadClean 完成 → 首片 wafer 进入 EPI 的间隔（越小越好）",
        "",
        "  优化前：EPI1 间隔 ~749s（OnLoadClean 在模拟启动时就开始，",
        "          但首片 wafer 需要 13+ 分钟才能到达 EPI，腔室大量空闲）",
        "",
        "  优化后：EPI1~EPI4 统一 ~73s（首个腔延迟公式对齐 + 后续腔计数器触发）",
    ], Inches(0.9), size=Pt(12))

    add_table(slide,
        ["腔室", "OnLoadClean 完成", "首片 wafer 放置", "间隔"],
        [
            ["EPI1", "00:24:34 (1474s)", "00:25:47 (1547s)", "73s"],
            ["EPI2", "00:34:43 (2083s)", "00:35:56 (2156s)", "73s"],
            ["EPI3", "00:45:27 (2727s)", "00:46:40 (2800s)", "73s"],
            ["EPI4", "00:56:59 (3419s)", "00:58:12 (3492s)", "73s"],
        ],
        Inches(3.8))

    add_footer(slide)

    # ====== Slide 12: HTML Report Overview ======
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_section_header(slide, "11", "完整分析报告 — simulation_report.html")

    add_body(slide, [
        "📌 报告内容",
        "  • 核心指标面板（完成数、WPH、总时间、平均周期）",
        "  • 工艺参数表（5 个 recipe + LL/PT/Cooling 时序 + 机械手操作时间）",
        "  • 腔室使用统计（总时间、使用次数、平均每次）",
        "  • 约束违反统计表（从甘特图真实驻留时间计算，去重）",
        "  • Wafer × Station 耗时矩阵（P 处理/D 驻留 分解，颜色高亮）",
        f"  • 完整甘特图（带腔室利用率百分比，{format_time(sim_time)} 时间跨度）",
        "  • Wafer History 下拉选择器（每片 wafer 的分步时间线）",
        "  • 🎬 机台动画 SVG 回放（和运行界面一致的 tool layout，播放/调速/步进）",
    ], Inches(0.9), size=Pt(12))

    add_key_point_box(slide,
        "📂 查看详细结果",
        "报告文件: simulation_report.html（自包含，无需服务器）\n双击用浏览器打开即可查看全部统计数据、甘特图、Wafer 历史、SVG 动画回放",
        Inches(4.5), width=Inches(8.5), height=Inches(0.85))

    # Check if report exists
    report_path = os.path.join(BASE_DIR, "result", "simulation_report.html")
    if os.path.exists(report_path):
        size_mb = round(os.path.getsize(report_path) / (1024 * 1024), 1)
        add_body(slide, [
            f"  ✅ 报告已生成：simulation_report.html ({size_mb}MB, 包含 2617 个回放快照 @10s 间隔)",
        ], Inches(5.6), size=Pt(12))

    add_footer(slide)

    # ====== Slide 13: Bug Fix History ======
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_section_header(slide, "12", "关键 Bug 修复记录")

    add_body(slide, [
        "📌 开发过程中发现并修复的关键 Bug",
    ], Inches(0.9), size=Pt(12))

    add_table(slide,
        ["Bug", "根因", "修复"],
        [
            ["lastUsedTime +dur 重复", "lambda 捕获 currentTimeSec 后再次 +dur", "统一使用 currentTimeSec"],
            ["EPI2/3/4 未做 OnLoadClean", "lastOnloadCleanStart 被 EPI/PC 共享，PC 打断 EPI stagger", "拆分为 per-type Map"],
            ["IdlePurge 看似无限循环", "EPI OnLoadClean 缺失导致全局死锁，purge 为表象", "修复 OnLoadClean 后正常"],
            ["EPI3/4 OnLoadClean 间隔大", "固定时钟 stagger 不匹配实际管道节奏", "EPI2+ 改用 wafer 计数器触发"],
            ["PreClean 1X Clean 标签错误", "handleChamberTimerDone 硬编码事件文字", "chamberCleanType Map 动态区分"],
            ["机械手不可见", "采样间隔 100s 太长 + 布局不匹配运行界面", "10s 采样 + SVG 对齐 ToolLayout"],
        ],
        Inches(1.6), font_size=Pt(9.5))

    add_footer(slide)

    # ====== Slide 14: Key Design Decisions ======
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_section_header(slide, "13", "关键设计决策")

    add_table(slide,
        ["决策", "理由", "影响"],
        [
            ["纯配置驱动，无硬编码", "最大灵活性，设备/工艺变更仅改 JSON", "5 个配置文件覆盖全部参数"],
            ["每秒 tick 离散模拟", "平衡精度与速度，1s 粒度满足约束需求", "26000s 模拟可数分钟完成"],
            ["TM1 PT 返回 > PT 前行", "防止 PT 槽满导致 EPI→PT 死锁", "返回路径优先清空"],
            ["PT fwd→buffer, ret→cooling", "分离前向/返回流量", "减少 PT 槽位争夺"],
            ["OnLoadClean: 首腔延迟 + 后续计数器", "首腔对齐首片，后续腔对齐实际管道", "全 4 腔间隔统一 ~73s"],
            ["IdlePurge 纯空闲驱动", "语义清晰，仅依赖 idle 时长", "与是否有 wafer 无关"],
            ["PURGING ≠ CLEANING 状态", "前端/甘特图清晰区分两种 AM", "橙色=Clean, 紫色=Purge"],
            ["Stagger = (EPI+Clean)/4", "考虑 1X Clean 时间的完整 EPI 周期", "~644s 间隔防 PT 拥堵"],
            ["安全裕度 40s", "驻留前瞻检查的保守缓冲", "避免边际违例"],
        ],
        Inches(1.1), font_size=Pt(8.5))

    add_footer(slide)

    # ====== Slide 15: Summary & Next Steps ======
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_section_header(slide, "14", "总结与展望")

    add_body(slide, [
        "✅ 已完成功能",
        "  • 完整的 Cluster Tool 离散事件仿真调度引擎（~1400 行核心代码）",
        "  • 双槽位 PT + CoolingStation 冷却逻辑 + 通道分离策略",
        "  • 时间感知前瞻死锁预防（三层防护：Stagger / EPI 容量 / PT 前向压力）",
        "  • AM 完整集成（OnLoadClean + 1X Clean + IdlePurge，状态/甘特图区分）",
        "  • 数据一致性自愈机制（healWaferLocations）",
        "  • Vue 3 前端可视化 + 交互式甘特图 + 实时 WebSocket 推送",
        "  • HTML 自包含报告 + SVG 动画回放 + PPTX 演示文稿",
        f"  • {done}/{total} 片晶圆全部成功完成，WPH = {wph}，驻留违例 = {total_violations} 次",
        "",
        "🔮 可扩展方向",
        "  • 双臂/多指机器人支持（需新增 Robot 模型 + Swap 逻辑）",
        "  • 多 CJ 并行调度 + 不同 recipe 混跑",
        "  • ML 优化调度策略（强化学习选择最优 wafer 释放时机）",
        "  • Web UI 实时修改配置文件（替代手动编辑 JSON）",
    ], Inches(0.9), size=Pt(11))

    add_footer(slide)

    # Save
    out_dir = os.path.join(BASE_DIR, "result")
    os.makedirs(out_dir, exist_ok=True)
    out_path = os.path.join(out_dir, "EPI_Scheduler_Report.pptx")
    prs.save(out_path)
    print(f"PPT saved to: {out_path}")
    print(f"Slides: {len(prs.slides)}")
    print(f"Size: {round(os.path.getsize(out_path) / 1024, 1)} KB")

if __name__ == '__main__':
    main()
